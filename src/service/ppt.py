# -*- coding: utf-8 -*-
from pptx import Presentation


def makeppt(names,lyrics):
    outputfile = "/home/gracechen/Documents/LineBot/pythonppt.pptx"
    f = open('/home/gracechen/Documents/LineBot/src/service/template.pptx',"rb")
    prs = Presentation(f)
    title_slide_layout = prs.slide_layouts[1] 

    for i in range(len(names)):
        name = names[i]
        lyric = lyrics[i]
        
        pages = lyric.split("\n\n")
        for l in pages:
            slide = prs.slides.add_slide(title_slide_layout)
            content = slide.placeholders[1]
            content.text = l
            title = slide.placeholders[0]
            title.text = name
    
    prs.save(outputfile)
    f.close()
    return outputfile
    



if __name__ == '__main__':
    name = ["我要歌頌你的力量"]
    lyric = ['''你救我脫離死亡 你救我脫離驚惶
在敵人面前 你為我擺設筵席
你恩典隨時夠用 你慈愛無盡無窮
在困難中 我要這樣說

就算是四面受敵 就算是暴風暴雨
當我要倒下 你卻伸手扶持我
你已將莫大能力 顯明在我的心裡
藉著信心 勇敢向前走

所以我要歌頌你的力量
早晨我要高唱你的慈愛
因為你作過我的高台
在急難中作過我的避難所

所以我要歌頌你的力量
早晨我要高唱你的慈愛
因為你是我的高台
是賜恩於我的神''']
    makeppt(name,lyric)
    
